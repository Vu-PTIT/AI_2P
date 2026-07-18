"""RAG retrieval backed by SentenceTransformers/FPT API and Qdrant."""

import json
from pathlib import Path
import os
import re
from typing import List, Optional
from uuid import uuid4

from model_errors import ModelUnavailableError


class DocumentChunk:
    """A chunk of source material."""

    def __init__(self, text: str, source: str, metadata: Optional[dict] = None):
        self.text = text
        self.source = source
        self.metadata = metadata or {}


class RAGEngine:
    """Embed meeting documents and retrieve relevant chunks from Qdrant."""

    def __init__(
        self,
        encoder: object | None = None,
        db: object | None = None,
        collection_name: str | None = None,
    ):
        self.embedding_model = os.getenv("RAG_EMBEDDING_MODEL", "sentence-transformers/all-MiniLM-L6-v2")
        self.db_path = os.getenv(
            "RAG_QDRANT_PATH",
            str(Path(__file__).resolve().parents[1] / "data" / "qdrant_local"),
        )
        self.collection_name = collection_name or os.getenv("RAG_COLLECTION", "meeting_docs")
        self._encoder = encoder
        self._db = db
        self._models = None
        # FPT AI Marketplace integration
        self._use_fpt_embedding = os.getenv("FPT_EMBEDDING", "").lower() == "true"
        self._use_fpt_reranker = os.getenv("FPT_RERANKER", "").lower() == "true"
        self._fpt_client = None

    def ingest_document(self, file_path: str) -> int:
        path = Path(file_path)
        chunks = self._chunk_text(self._read_document(path))
        self._upsert_chunks(
            [
                DocumentChunk(text=chunk, source=str(path), metadata={"chunk": index})
                for index, chunk in enumerate(chunks)
            ],
        )
        return len(chunks)

    def retrieve(self, query: str, top_k: int = 5) -> List[DocumentChunk]:
        if not query.strip():
            return []

        encoder, db = self._ensure_backend()
        vector = self._encode_text(encoder, query)
        # ponytail: fetch more candidates when reranker is active, rerank trims to top_k
        fetch_k = top_k * 3 if self._use_fpt_reranker else top_k
        try:
            results = db.search(
                collection_name=self.collection_name,
                query_vector=vector,
                limit=fetch_k,
            )
        except AttributeError:
            response = db.query_points(
                collection_name=self.collection_name,
                query=vector,
                limit=fetch_k,
            )
            results = response.points

        chunks: List[DocumentChunk] = []
        for result in results:
            payload = result.payload if hasattr(result, "payload") else result.get("payload", {})
            chunks.append(
                DocumentChunk(
                    text=payload.get("text", ""),
                    source=payload.get("source", ""),
                    metadata=payload.get("metadata", {}),
                ),
            )
        chunks = [chunk for chunk in chunks if chunk.text]

        if self._use_fpt_reranker and chunks:
            chunks = self._fpt_rerank(query, chunks, top_k)

        return chunks[:top_k]

    def add_session_transcript(self, transcript: str, session_id: str) -> None:
        chunks = [
            DocumentChunk(
                text=chunk,
                source=f"session:{session_id}",
                metadata={"chunk": index, "session_id": session_id},
            )
            for index, chunk in enumerate(self._chunk_text(transcript))
        ]
        self._upsert_chunks(chunks)

    def _ensure_backend(self):
        if self._encoder is not None and self._db is not None:
            return self._encoder, self._db

        try:
            from qdrant_client import QdrantClient, models
        except ImportError as error:
            raise ModelUnavailableError(
                "qdrant-client is required for RAG.",
            ) from error

        try:
            self._models = models
            if self._use_fpt_embedding:
                # ponytail: FPT API embedding — encoder is this engine itself
                self._encoder = self
            else:
                from sentence_transformers import SentenceTransformer
                self._encoder = SentenceTransformer(self.embedding_model)
            self._db = QdrantClient(path=self.db_path)
            self._ensure_collection()
        except Exception as error:
            raise ModelUnavailableError(
                f"RAG backend is unavailable. Set RAG_EMBEDDING_MODEL/RAG_QDRANT_PATH correctly.",
            ) from error

        return self._encoder, self._db

    def _ensure_collection(self) -> None:
        if self._models is None:
            return

        try:
            self._db.get_collection(self.collection_name)
            return
        except Exception:
            pass

        dimension = self._embedding_dimension()
        self._db.create_collection(
            collection_name=self.collection_name,
            vectors_config=self._models.VectorParams(
                size=dimension,
                distance=self._models.Distance.COSINE,
            ),
        )

    def _embedding_dimension(self) -> int:
        if hasattr(self._encoder, "get_sentence_embedding_dimension"):
            dimension = self._encoder.get_sentence_embedding_dimension()
            if dimension:
                return int(dimension)
        return len(self._as_vector(self._encoder.encode("dimension probe")))

    def _upsert_chunks(self, chunks: List[DocumentChunk]) -> None:
        if not chunks:
            return

        encoder, db = self._ensure_backend()
        texts = [chunk.text for chunk in chunks]
        vectors = self._encode_texts(encoder, texts)
        points = []
        for chunk, vector in zip(chunks, vectors):
            payload = {
                "text": chunk.text,
                "source": chunk.source,
                "metadata": chunk.metadata,
            }
            point_id = str(uuid4())
            point_vector = self._as_vector(vector)
            if self._models is not None:
                points.append(
                    self._models.PointStruct(
                        id=point_id,
                        vector=point_vector,
                        payload=payload,
                    ),
                )
            else:
                points.append({"id": point_id, "vector": point_vector, "payload": payload})
        db.upsert(collection_name=self.collection_name, points=points)

    def _chunk_text(self, text: str, max_chars: int = 800) -> List[str]:
        chunks: List[str] = []
        for paragraph in re.split(r"\n\s*\n", text):
            paragraph = paragraph.strip()
            if not paragraph:
                continue
            while len(paragraph) > max_chars:
                chunks.append(paragraph[:max_chars].strip())
                paragraph = paragraph[max_chars:].strip()
            chunks.append(paragraph)
        return chunks

    def _read_document(self, path: Path) -> str:
        suffix = path.suffix.casefold()
        if suffix in ("", ".txt", ".md", ".markdown", ".csv", ".tsv", ".log"):
            return path.read_text(encoding="utf-8-sig")
        if suffix == ".json":
            return json.dumps(
                json.loads(path.read_text(encoding="utf-8-sig")),
                ensure_ascii=False,
                indent=2,
            )
        if suffix == ".pdf":
            return self._read_pdf(path)
        if suffix == ".docx":
            return self._read_docx(path)
        if suffix == ".pptx":
            return self._read_pptx(path)
        raise ModelUnavailableError(f"Unsupported RAG document type: {suffix or '<none>'}")

    def _read_pdf(self, path: Path) -> str:
        try:
            from pypdf import PdfReader
        except ImportError as error:
            raise ModelUnavailableError("pypdf is required to ingest PDF documents.") from error
        reader = PdfReader(str(path))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)

    def _read_docx(self, path: Path) -> str:
        try:
            from docx import Document
        except ImportError as error:
            raise ModelUnavailableError("python-docx is required to ingest DOCX documents.") from error
        document = Document(str(path))
        return "\n\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)

    def _read_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError as error:
            raise ModelUnavailableError("python-pptx is required to ingest PPTX documents.") from error
        presentation = Presentation(str(path))
        lines = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(f"Slide {slide_index}: {shape.text}")
        return "\n\n".join(lines)

    def _as_vector(self, value) -> list[float]:
        return value.tolist() if hasattr(value, "tolist") else list(value)

    # --- FPT AI Marketplace helpers ---

    def _get_fpt_client(self):
        if self._fpt_client is None:
            from config.fpt_models import get_fpt_client
            self._fpt_client = get_fpt_client()
        return self._fpt_client

    def encode(self, text):
        """SentenceTransformer-compatible encode() via FPT API."""
        return self._fpt_embed([text])[0]

    def _encode_text(self, encoder, text: str) -> list[float]:
        if self._use_fpt_embedding:
            return self._fpt_embed([text])[0]
        return self._as_vector(encoder.encode(text))

    def _encode_texts(self, encoder, texts: list[str]) -> list:
        if self._use_fpt_embedding:
            return self._fpt_embed(texts)
        return encoder.encode(texts)

    def _fpt_embed(self, texts: list[str]) -> list[list[float]]:
        """Call FPT /embeddings endpoint with Vietnamese_Embedding."""
        from config.fpt_models import FPT_BASE_URL, FPT_API_KEY, MODELS as FPT
        import requests

        response = requests.post(
            f"{FPT_BASE_URL.rstrip('/')}/embeddings",
            headers={
                "Authorization": f"Bearer {FPT_API_KEY}",
                "Content-Type": "application/json",
            },
            json={"input": texts, "model": FPT["embedding"]},
            timeout=10,
        )
        response.raise_for_status()
        data = response.json()
        # Sort by index to preserve order
        items = sorted(data.get("data", []), key=lambda x: x.get("index", 0))
        return [item["embedding"] for item in items]

    def _fpt_rerank(self, query: str, chunks: List[DocumentChunk], top_n: int) -> List[DocumentChunk]:
        """Call FPT /v1/rerank with bge-reranker-v2-m3 to re-order chunks."""
        from config.fpt_models import FPT_BASE_URL, FPT_API_KEY, MODELS as FPT
        import requests

        try:
            response = requests.post(
                f"{FPT_BASE_URL.rstrip('/')}/v1/rerank",
                headers={
                    "Authorization": f"Bearer {FPT_API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": FPT["reranker"],
                    "query": query,
                    "documents": [chunk.text for chunk in chunks],
                    "top_n": top_n,
                },
                timeout=5,
            )
            response.raise_for_status()
            results = response.json().get("results", [])
            return [chunks[item["index"]] for item in results if item.get("index", -1) < len(chunks)]
        except Exception:
            # ponytail: reranker failure is non-fatal, return original order
            return chunks[:top_n]
